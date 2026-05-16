from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

# Color palette
PRIMARY = RGBColor(0x1E, 0x40, 0xAF)    # Deep blue
SECONDARY = RGBColor(0x7C, 0x3A, 0xED)  # Purple
ACCENT = RGBColor(0x06, 0xB6, 0xD4)     # Cyan
DARK = RGBColor(0x1E, 0x29, 0x3B)       # Dark gray
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)   # Light gray bg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY = RGBColor(0x64, 0x74, 0x8B)
SUCCESS = RGBColor(0x10, 0xB9, 0x81)     # Green
WARNING = RGBColor(0xF5, 0x9E, 0x0B)     # Amber

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_after = Pt(8)
    return txBox


# ============================================================
# Slide 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK)

# Accent bar at top
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT)

# Subtitle
add_text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(0.6),
         "CODEBASE OVERVIEW", size=18, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

# Title
add_text(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
         "BlogFlow", size=54, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Description
add_text(slide, Inches(1.5), Inches(4.2), Inches(8), Inches(1),
         "AI-Powered Blog Content Creation Platform", size=24, color=RGBColor(0x94, 0xA3, 0xB8))

# Tech tags
tags = ["React 18", "TypeScript", "Supabase", "Claude AI", "Tailwind CSS v4", "Vite"]
x_start = Inches(1.5)
for tag in tags:
    shape = add_shape_bg(slide, x_start, Inches(5.5), Inches(1.6), Inches(0.45), RGBColor(0x33, 0x41, 0x55), radius=True)
    shape.text_frame.paragraphs[0].text = tag
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    shape.text_frame.paragraphs[0].font.name = 'Arial'
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    x_start += Inches(1.75)


# ============================================================
# Slide 2: Project Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
         "Project Overview", size=32, bold=True, color=DARK)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
         "BlogFlow는 AI 기반 블로그 콘텐츠 제작 플랫폼입니다.\n"
         "사용자가 URL/파일을 입력하면, AI가 인사이트를 추출하고 리서치를 수행한 뒤\n"
         "아웃라인 작성과 최종 드래프트 생성, SEO 최적화까지 자동화합니다.",
         size=16, color=TEXT_GRAY)

# Feature cards
features = [
    ("Resource Input", "URL/파일에서\n콘텐츠 수집", PRIMARY),
    ("AI Analysis", "Claude AI로\n인사이트 추출", SECONDARY),
    ("Deep Research", "심층 리서치 및\n시장 분석", ACCENT),
    ("Outline Editor", "3-패널 에디터로\n구조 편집", SUCCESS),
    ("Draft Generation", "AI 기반\n드래프트 작성", WARNING),
    ("SEO Analysis", "SEO 점수 분석\n및 최적화", RGBColor(0xEF, 0x44, 0x44)),
]

for i, (title, desc, color) in enumerate(features):
    x = Inches(0.8 + i * 2.05)
    y = Inches(2.8)

    card = add_shape_bg(slide, x, y, Inches(1.85), Inches(3.5), LIGHT_BG, radius=True)
    # Color bar at top of card
    bar = add_shape_bg(slide, x + Inches(0.15), y + Inches(0.2), Inches(1.55), Inches(0.06), color)

    add_text(slide, x + Inches(0.15), y + Inches(0.5), Inches(1.55), Inches(0.5),
             title, size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.15), y + Inches(1.2), Inches(1.55), Inches(1.5),
             desc, size=13, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Step number
    add_text(slide, x + Inches(0.15), y + Inches(2.8), Inches(1.55), Inches(0.4),
             f"Step {i+1}", size=11, bold=True, color=color, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 3: Workflow State Machine
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
         "Workflow State Machine", size=32, bold=True, color=DARK)

add_text(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5),
         "7단계 순차 워크플로우 — useWorkflow() 커스텀 훅으로 상태 관리",
         size=16, color=TEXT_GRAY)

steps = [
    ("input", "ResourceInput", "URL/파일 입력\n키워드, 타겟 설정", PRIMARY),
    ("analyzing", "AnalysisLoading", "리소스 수집\n인사이트 추출", RGBColor(0x63, 0x66, 0xF1)),
    ("selection", "InsightSelection", "인사이트 선택\n개발할 주제 결정", SECONDARY),
    ("researching", "DeepResearch", "심층 리서치\n시장/경쟁 분석", ACCENT),
    ("outline", "OutlineEditor", "3-패널 에디터\n구조 편집", SUCCESS),
    ("writing", "WritingDraft", "AI 드래프트\n생성 중", WARNING),
    ("final", "FinalDraftScreen", "SEO 분석\n발행/내보내기", RGBColor(0xEF, 0x44, 0x44)),
]

for i, (step, component, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 1.78)
    y = Inches(2.3)

    # Step box
    box = add_shape_bg(slide, x, y, Inches(1.58), Inches(3.8), LIGHT_BG, radius=True)

    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.54), y + Inches(0.2), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = str(i + 1)
    circle.text_frame.paragraphs[0].font.size = Pt(16)
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Step name
    add_text(slide, x + Inches(0.05), y + Inches(0.85), Inches(1.48), Inches(0.35),
             f"'{step}'", size=12, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Component name
    add_text(slide, x + Inches(0.05), y + Inches(1.3), Inches(1.48), Inches(0.35),
             component, size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Description
    add_text(slide, x + Inches(0.05), y + Inches(1.9), Inches(1.48), Inches(1),
             desc, size=11, color=DARK, align=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        x + Inches(1.58), y + Inches(1.6),
                                        Inches(0.2), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        arrow.line.fill.background()


# ============================================================
# Slide 4: Architecture Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
         "Architecture Overview", size=32, bold=True, color=DARK)

# Frontend layer
add_shape_bg(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), RGBColor(0xEF, 0xF6, 0xFF), radius=True)
add_text(slide, Inches(0.7), Inches(1.6), Inches(2), Inches(0.4),
         "Frontend (React + Vite)", size=14, bold=True, color=PRIMARY)

fe_items = [
    ("App.tsx", "라우팅 (/, /saved, /drafts/:id)"),
    ("WorkflowContainer", "워크플로우 단계 렌더링"),
    ("useWorkflow()", "상태 머신 훅"),
    ("blogApi", "API 레이어 (~540줄)"),
    ("components/", "shadcn/ui + 공유 컴포넌트"),
]
for j, (name, desc) in enumerate(fe_items):
    yy = Inches(2.2 + j * 0.85)
    card = add_shape_bg(slide, Inches(0.8), yy, Inches(5.2), Inches(0.7), WHITE, radius=True)
    add_text(slide, Inches(1.0), yy + Inches(0.05), Inches(2), Inches(0.3),
             name, size=13, bold=True, color=DARK)
    add_text(slide, Inches(1.0), yy + Inches(0.35), Inches(4.5), Inches(0.3),
             desc, size=11, color=TEXT_GRAY)

# Backend layer
add_shape_bg(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.8), RGBColor(0xF0, 0xFD, 0xF4), radius=True)
add_text(slide, Inches(7.0), Inches(1.6), Inches(3), Inches(0.4),
         "Backend (Supabase Edge Functions)", size=14, bold=True, color=SUCCESS)

be_items = [
    "collect-resource — URL 스크래핑",
    "extract-insights — Claude AI 분석",
    "deep-research — 심층 리서치",
    "generate-outline — 아웃라인 생성",
    "write-draft — 드래프트 작성",
    "analyze-seo — SEO 분석",
]
for j, item in enumerate(be_items):
    add_text(slide, Inches(7.2), Inches(2.2 + j * 0.32), Inches(5), Inches(0.3),
             f"• {item}", size=11, color=DARK)

# Database layer
add_shape_bg(slide, Inches(6.8), Inches(4.6), Inches(6), Inches(2.4), RGBColor(0xFE, 0xF3, 0xC7), radius=True)
add_text(slide, Inches(7.0), Inches(4.7), Inches(3), Inches(0.4),
         "Database (Supabase PostgreSQL)", size=14, bold=True, color=WARNING)

db_items = [
    "workflow_sessions — 워크플로우 세션",
    "resources — 수집된 리소스",
    "insights — 추출된 인사이트",
    "research — 심층 리서치 결과",
    "outlines / drafts — 아웃라인 & 드래프트",
]
for j, item in enumerate(db_items):
    add_text(slide, Inches(7.2), Inches(5.3 + j * 0.3), Inches(5), Inches(0.3),
             f"• {item}", size=11, color=DARK)


# ============================================================
# Slide 5: Directory Structure
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
         "Directory Structure", size=32, bold=True, color=DARK)

# Left column - src/
add_shape_bg(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(5.6), LIGHT_BG, radius=True)
add_text(slide, Inches(0.7), Inches(1.5), Inches(2), Inches(0.4),
         "src/", size=18, bold=True, color=PRIMARY)

src_tree = """App.tsx                  — 메인 라우터
main.tsx                 — 엔트리 포인트
components/shared/       — Layout, Sidebar, ProgressBar
components/ui/           — shadcn/ui 컴포넌트
features/workflow/       — 워크플로우 (핵심 기능)
  ├ WorkflowContainer.tsx — 단계별 렌더링
  ├ hooks/useWorkflow.ts  — 상태 머신 훅
  ├ components/           — 6개 단계 컴포넌트
  └ types/index.ts        — TypeScript 타입 정의
features/drafts/         — 저장된 드래프트 관리
features/prompts/        — AI 프롬프트 관리 (beta)
features/news/           — 뉴스 검색
lib/api.ts               — Supabase API (~540줄)
lib/supabase.ts          — Supabase 클라이언트"""

add_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
         src_tree, size=11, color=DARK, font_name='Consolas')

# Right column - supabase/
add_shape_bg(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(5.6), LIGHT_BG, radius=True)
add_text(slide, Inches(7.0), Inches(1.5), Inches(2), Inches(0.4),
         "supabase/", size=18, bold=True, color=SUCCESS)

supa_tree = """migrations/
  ├ 001_initial_schema.sql
  ├ 002_add_seo_fields.sql
  ├ 003_add_prompts_table.sql
  ├ 004_seed_prompts.sql
  └ 005-008 추가 마이그레이션

functions/
  ├ collect-resource/     — URL 스크래핑
  ├ extract-insights/     — Claude AI 분석
  ├ deep-research/        — 심층 리서치
  ├ generate-outline/     — 아웃라인 생성
  ├ write-draft/          — 드래프트 작성
  ├ analyze-seo/          — SEO 분석
  ├ search-news/          — 뉴스 검색
  └ _shared/              — 공유 유틸리티
      ├ anthropic.ts
      ├ supabase.ts
      ├ cors.ts
      └ rateLimit.ts"""

add_text(slide, Inches(7.1), Inches(2.0), Inches(5.5), Inches(4.5),
         supa_tree, size=11, color=DARK, font_name='Consolas')


# ============================================================
# Slide 6: Database Schema
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
         "Database Schema", size=32, bold=True, color=DARK)

add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
         "Supabase PostgreSQL — 6개 핵심 테이블 + JSONB 유연한 데이터 구조",
         size=14, color=TEXT_GRAY)

# Table cards
tables = [
    ("workflow_sessions", "세션 관리", ["id (UUID PK)", "status (워크플로우 단계)", "keywords", "target_audience"], PRIMARY),
    ("resources", "수집 리소스", ["session_id (FK)", "source_type (url|file)", "title, content", "collected_at"], ACCENT),
    ("insights", "추출 인사이트", ["title, signal", "confidence (high|mid|low)", "tags TEXT[]", "status (pending|selected)"], SECONDARY),
    ("research", "심층 리서치", ["insight_id (FK)", "market_data JSONB", "competitor_analysis", "statistics, sources JSONB"], SUCCESS),
    ("outlines", "아웃라인", ["title, thesis, tone", "sections JSONB", "structure_pattern", "status (draft|approved)"], WARNING),
    ("drafts", "최종 드래프트", ["content (markdown)", "word_count, char_count", "seo_metrics JSONB", "status (draft|final|pub)"], RGBColor(0xEF, 0x44, 0x44)),
]

for i, (name, label, fields, color) in enumerate(tables):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.8 + row * 2.8)

    card = add_shape_bg(slide, x, y, Inches(3.9), Inches(2.5), LIGHT_BG, radius=True)
    # Header bar
    add_shape_bg(slide, x, y, Inches(3.9), Inches(0.5), color, radius=True)
    add_text(slide, x + Inches(0.15), y + Inches(0.05), Inches(3.5), Inches(0.4),
             f"{name}  —  {label}", size=13, bold=True, color=WHITE)

    for j, field in enumerate(fields):
        add_text(slide, x + Inches(0.2), y + Inches(0.65 + j * 0.4), Inches(3.5), Inches(0.35),
                 f"• {field}", size=11, color=DARK)


# ============================================================
# Slide 7: API Layer
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
         "API Layer (blogApi)", size=32, bold=True, color=DARK)

add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
         "src/lib/api.ts — 중앙집중식 Supabase API 래퍼 (~540줄)",
         size=14, color=TEXT_GRAY)

api_groups = [
    ("Session", ["createSession()", "getSession()", "updateSessionStatus()"], PRIMARY),
    ("Resources", ["collectResource()", "addFileResource()", "getResources()"], ACCENT),
    ("Insights", ["extractInsights()", "getInsights()", "selectInsights()"], SECONDARY),
    ("Research", ["deepResearch()", "getResearch()"], SUCCESS),
    ("Outline", ["generateOutline()", "getOutline()", "updateOutline()", "approveOutline()"], WARNING),
    ("Draft", ["writeDraft()", "getDraft()", "updateDraft()", "publishDraft()", "getAllDrafts()", "deleteDraft()"], RGBColor(0xEF, 0x44, 0x44)),
    ("SEO", ["analyzeSeo()", "getSeoMetrics()"], RGBColor(0x8B, 0x5C, 0xF6)),
    ("Prompts", ["getAllPrompts()", "updatePrompt()"], RGBColor(0xEC, 0x48, 0x99)),
]

for i, (group, methods, color) in enumerate(api_groups):
    col = i % 4
    row = i // 4
    x = Inches(0.5 + col * 3.15)
    y = Inches(1.8 + row * 2.8)

    card = add_shape_bg(slide, x, y, Inches(2.95), Inches(2.5), LIGHT_BG, radius=True)
    add_shape_bg(slide, x, y, Inches(2.95), Inches(0.45), color, radius=True)
    add_text(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.5), Inches(0.35),
             group, size=14, bold=True, color=WHITE)

    for j, method in enumerate(methods):
        add_text(slide, x + Inches(0.15), y + Inches(0.55 + j * 0.3), Inches(2.6), Inches(0.28),
                 method, size=11, color=DARK, font_name='Consolas')


# ============================================================
# Slide 8: Key Design Patterns
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
         "Key Design Patterns", size=32, bold=True, color=DARK)

patterns = [
    ("State Machine Hook",
     "useWorkflow() 커스텀 훅이 전체 워크플로우 상태와\n작업을 캡슐화. 단일 진실 원천(Single Source of Truth)",
     PRIMARY),
    ("Component-Driven Workflow",
     "각 워크플로우 단계가 독립적인 컴포넌트로 분리.\nWorkflowContainer에서 조건부 렌더링",
     SECONDARY),
    ("3-Panel Layout (Outline Editor)",
     "좌측: 아웃라인 트리 / 중앙: 상세 편집\n우측: AI 어시스턴트 — 전문적 UX 패턴",
     ACCENT),
    ("Centralized API Layer",
     "모든 Supabase 호출이 blogApi 객체를 통해 이루어짐.\n일관된 에러 처리와 타입 안전성 보장",
     SUCCESS),
    ("Server-side AI Processing",
     "Claude API 호출은 Edge Function에서 실행.\nAPI 키 보호 + 서버사이드 처리",
     WARNING),
    ("JSONB for Flexible Data",
     "research, outlines, drafts에 JSONB 활용.\n스키마 변경 없이 유연한 데이터 구조",
     RGBColor(0xEF, 0x44, 0x44)),
]

for i, (title, desc, color) in enumerate(patterns):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.5 + row * 2.8)

    card = add_shape_bg(slide, x, y, Inches(3.95), Inches(2.5), LIGHT_BG, radius=True)
    bar = add_shape_bg(slide, x + Inches(0.15), y + Inches(0.2), Inches(0.08), Inches(2.1), color)

    add_text(slide, x + Inches(0.4), y + Inches(0.2), Inches(3.3), Inches(0.5),
             title, size=16, bold=True, color=DARK)

    add_text(slide, x + Inches(0.4), y + Inches(0.8), Inches(3.3), Inches(1.5),
             desc, size=12, color=TEXT_GRAY)


# ============================================================
# Slide 9: Tech Stack Detail
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
         "Tech Stack & Dependencies", size=32, bold=True, color=DARK)

stacks = [
    ("Core Framework", [
        "React 18 — UI 라이브러리",
        "TypeScript — 타입 안전성",
        "Vite — 빌드 툴 & 개발 서버",
        "React Router v6 — 클라이언트 라우팅",
    ], PRIMARY),
    ("Styling & UI", [
        "Tailwind CSS v4 — 유틸리티 CSS",
        "shadcn/ui — Radix UI 컴포넌트",
        "Lucide React — 아이콘",
        "Motion (Framer) — 애니메이션",
    ], SECONDARY),
    ("Backend & Data", [
        "Supabase — PostgreSQL + Auth",
        "Edge Functions (Deno) — 서버리스",
        "Claude API — AI 처리",
        "Perplexity API — 뉴스 검색",
    ], SUCCESS),
    ("Dev Tools", [
        "React Hook Form — 폼 관리",
        "Path alias @ → ./src",
        "ESLint + TypeScript config",
        "npm 패키지 매니저",
    ], WARNING),
]

for i, (category, items, color) in enumerate(stacks):
    x = Inches(0.5 + i * 3.15)
    y = Inches(1.5)

    card = add_shape_bg(slide, x, y, Inches(2.95), Inches(5), LIGHT_BG, radius=True)
    add_shape_bg(slide, x, y, Inches(2.95), Inches(0.5), color, radius=True)
    add_text(slide, x + Inches(0.15), y + Inches(0.07), Inches(2.5), Inches(0.4),
             category, size=15, bold=True, color=WHITE)

    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.2), y + Inches(0.7 + j * 0.9), Inches(2.5), Inches(0.8),
                 item, size=12, color=DARK)


# ============================================================
# Slide 10: Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT)

add_text(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(0.7),
         "Summary", size=36, bold=True, color=WHITE)

summary_items = [
    "7단계 워크플로우 상태 머신으로 명확한 프로세스 관리",
    "각 단계를 독립 컴포넌트로 분리한 깔끔한 아키텍처",
    "blogApi 객체를 통한 중앙집중식 API 관리 (~540줄)",
    "Edge Functions로 AI 호출을 서버사이드에서 안전하게 처리",
    "JSONB 활용으로 유연한 데이터 구조 & 스키마 확장성",
    "TypeScript 전체 적용으로 타입 안전성 보장",
    "shadcn/ui + Tailwind + Motion으로 세련된 UI/UX",
]

for i, item in enumerate(summary_items):
    bullet = add_text(slide, Inches(1.5), Inches(2.2 + i * 0.6), Inches(10), Inches(0.5),
                      f"  {item}", size=16, color=RGBColor(0xCB, 0xD5, 0xE1))

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(1.5), Inches(2.3 + i * 0.6),
                                     Inches(0.15), Inches(0.15))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()

add_text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
         "BlogFlow — AI-Powered Blog Content Creation Platform",
         size=14, color=TEXT_GRAY, align=PP_ALIGN.LEFT)


# Save
output_path = '/Users/awaji/Desktop/project/blog-web-agent/BlogFlow_Codebase_Overview.pptx'
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
